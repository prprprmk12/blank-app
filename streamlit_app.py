import streamlit as st
from mistralai import Mistral
from pptx import Presentation
from io import BytesIO
import re
import os
from dotenv import load_dotenv
import time
import urllib.parse
import http.client

# --- CONFIGURATION ---
load_dotenv() # Load variables from .env if present

MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY", "Ybw8mXxtjlIQIpy1xVSZU5Cap1V1unta")
POLLINATIONS_API_KEY = os.getenv("Pollinations_API_Key", "sk_HqwqD3mhFkpk7YDlFcKmg1qdfWEv6FUC") or os.getenv("POLLINATIONS_API_KEY")
PROJECT_NAME = "Платформа 2026: AI-Интеграция"

st.set_page_config(page_title=PROJECT_NAME, layout="wide", page_icon="🎓")

# --- UI STYLING ---
# --- UI STYLING ---
st.markdown("""
<style>
    /* Import Google Font */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap');

    html, body, [class*="css"]  {
        font-family: 'Inter', sans-serif;
    }

    /* Main App Background - Solid Premium Dark */
    .stApp {
        background-color: #0f172a; /* Slate 900 */
        background-image: 
            radial-gradient(at 0% 0%, rgba(56, 189, 248, 0.1) 0px, transparent 50%), 
            radial-gradient(at 100% 100%, rgba(139, 92, 246, 0.1) 0px, transparent 50%);
        background-attachment: fixed;
    }

    /* Sidebar Styling */
    section[data-testid="stSidebar"] {
        background-color: #1e293b; /* Slate 800 */
        border-right: 1px solid #334155;
    }
    section[data-testid="stSidebar"] h1, section[data-testid="stSidebar"] h2, section[data-testid="stSidebar"] h3 {
        color: #f8fafc !important;
    }
    section[data-testid="stSidebar"] label {
        color: #cbd5e1 !important;
    }

    /* Header Styling */
    .main-header {
        text-align: center;
        padding: 4rem 1rem;
        margin-bottom: 2rem;
    }
    .main-header h1 {
        font-weight: 800;
        font-size: 4rem;
        margin: 0;
        background: linear-gradient(to right, #38bdf8, #818cf8, #c084fc);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-shadow: 0 0 20px rgba(139, 92, 246, 0.3);
    }
    .main-header p {
        color: #94a3b8;
        font-size: 1.4rem;
        margin-top: 1rem;
        font-weight: 300;
    }

    /* Card/Container Styling - Dark Glass */
    .feature-card {
        background: rgba(30, 41, 59, 0.7);
        backdrop-filter: blur(12px);
        border-radius: 20px;
        padding: 30px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        box-shadow: 0 8px 32px 0 rgba(0, 0, 0, 0.3);
        color: #e2e8f0;
        margin-bottom: 20px;
    }
    
    /* Headings in cards/tabs */
    h1, h2, h3, h4, 
    .stMarkdown h1, .stMarkdown h2, .stMarkdown h3 {
        color: #f1f5f9 !important;
    }

    /* Text inputs */
    .stTextInput > div > div > input, 
    .stTextArea > div > div > textarea {
        background-color: #1e293b;
        color: #f8fafc;
        border: 1px solid #475569;
        border-radius: 12px;
    }
    .stTextInput > div > div > input:focus, 
    .stTextArea > div > div > textarea:focus {
        border-color: #818cf8;
        box-shadow: 0 0 0 3px rgba(129, 140, 248, 0.2);
    }

    /* Buttons */
    .stButton > button {
        background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        border-radius: 12px;
        font-weight: 600;
        letter-spacing: 0.5px;
        transition: all 0.3s ease;
        text-transform: uppercase;
        font-size: 0.9rem;
    }
    .stButton > button:hover {
        background: linear-gradient(135deg, #4338ca 0%, #6d28d9 100%);
        transform: translateY(-2px);
        box-shadow: 0 10px 20px -10px rgba(124, 58, 237, 0.5);
    }

    /* Tabs */
    .stTabs [data-baseweb="tab-list"] {
        background-color: #1e293b;
        padding: 8px;
        border-radius: 16px;
        border: 1px solid #334155;
    }
    .stTabs [data-baseweb="tab"] {
        color: #94a3b8;
        border-radius: 12px;
    }
    .stTabs [aria-selected="true"] {
        background-color: #334155;
        color: #ffffff;
    }
    
    /* Spinner */
    .stSpinner > div {
        border-top-color: #818cf8 !important;
    }
</style>
""", unsafe_allow_html=True)

st.markdown(f'''
<div class="main-header">
    <h1>{PROJECT_NAME}</h1>
    <p>Образовательная платформа нового поколения</p>
</div>
''', unsafe_allow_html=True)

# --- SIDEBAR: API SETUP ---
with st.sidebar:
    st.header("🔑 Настройки API")
    
    if MISTRAL_API_KEY:
        st.success("Mistral API Key загружен.")
    else:
        st.warning("Mistral API Key не найден в .env")
    
    st.info("ℹ️ Генерация изображений (gen.pollinations.ai)")
    pollinations_key = st.text_input("Pollinations API Key", value=POLLINATIONS_API_KEY or "", type="password", help="Получите ключ на enter.pollinations.ai. Ключ обязателен для этого API.")
    if not pollinations_key and not POLLINATIONS_API_KEY:
        st.warning("⚠️ Для генерации изображений требуется API ключ.")
    
# --- HELPER FUNCTIONS ---

def ask_mistral(prompt, system_prompt="Ты экспертный ИИ-ассистент образовательной платформы 2026."):
    if not MISTRAL_API_KEY: 
        return "Ошибка: Mistral API Key не настроен."
    try:
        client = Mistral(api_key=MISTRAL_API_KEY)
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Ошибка Mistral: {e}"

def generate_image_pollinations(prompt):
    """
    Generates an image by fetching bytes via http.client from gen.pollinations.ai.
    Requires an API key.
    """
    key_to_use = pollinations_key or POLLINATIONS_API_KEY
    if not key_to_use:
        st.error("Ошибка: API ключ Pollinations не настроен. Получите его на enter.pollinations.ai")
        return None

    encoded_prompt = urllib.parse.quote(prompt)
    seed = int(time.time())
    
    # Path with query parameters. Added model=flux for better results as seen in docs.
    path = f"/image/{encoded_prompt}?seed={seed}&width=1024&height=768&nologo=true&model=flux"
    
    # Setup headers - Bearer token is the standard
    headers = {
        "Authorization": f"Bearer {key_to_use}",
        "Content-Type": "application/json"
    }
    
    try:
        conn = http.client.HTTPSConnection("gen.pollinations.ai")
        conn.request("GET", path, headers=headers)
        response = conn.getresponse()
        
        if response.status == 200:
            image_bytes = response.read()
            conn.close()
            return image_bytes
        elif response.status == 401:
            conn.close()
            st.error("Ошибка 401: Неверный или отсутствующий API ключ. Проверьте настройки в боковой панели.")
            return None
        else:
            error_msg = response.read().decode()
            conn.close()
            st.error(f"Ошибка Pollinations ({response.status}): {error_msg}")
            return None
    except Exception as e:
        st.error(f"Ошибка соединения: {e}")
        return None

def create_pptx(content):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI Generated Presentation"
    slide.placeholders[1].text = content
    
    ppt_io = BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()

# --- TABS INTERFACE ---
tabs = st.tabs([
    "🎓 Репетитор", "💻 Ошибки", "📝 Тесты", "📖 Уроки", 
    "📊 Презентации", "✍️ Эссе", "🌐 Переводчик", "🖼 Иллюстратор"
])

# 1. AI-Repetitor
with tabs[0]:
    st.subheader("👨‍🏫 AI-Репетитор")
    user_query = st.text_area("Задай вопрос по любому предмету:", height=150)
    if st.button("Спросить"):
        if user_query:
            with st.spinner("Думаю..."):
                answer = ask_mistral(user_query)
                st.markdown(f'<div class="feature-card">{answer}</div>', unsafe_allow_html=True)

# 2. Error Explanation
with tabs[1]:
    st.subheader("🔍 Объяснение ошибок")
    code_input = st.text_area("Вставь код с ошибкой:", height=150)
    if st.button("Анализировать"):
        if code_input:
            with st.spinner("Разбираю код..."):
                explanation = ask_mistral(f"Найди ошибку в этом коде и объясни, почему она возникла:\n{code_input}")
                st.markdown(f'<div class="feature-card">{explanation}</div>', unsafe_allow_html=True)

# 3. Test Generator
with tabs[2]:
    st.subheader("📋 Генератор тестов")
    topic = st.text_input("Тема теста (например: Фотосинтез):")
    if st.button("Создать тест"):
        if topic:
            with st.spinner("Генерирую вопросы..."):
                test = ask_mistral(f"Создай тест из 5 вопросов с вариантами ответов на тему: {topic}")
                st.markdown(f'<div class="feature-card">{test}</div>', unsafe_allow_html=True)

# 4. Lesson Generator
with tabs[3]:
    st.subheader("📚 Планировщик уроков")
    lesson_topic = st.text_input("Тема урока:")
    if st.button("План урока"):
        if lesson_topic:
            with st.spinner("Составляю план..."):
                lesson = ask_mistral(f"Составь подробный план урока на 45 минут для темы: {lesson_topic}")
                st.markdown(f'<div class="feature-card">{lesson}</div>', unsafe_allow_html=True)

# 5. Presentation Generator
with tabs[4]:
    st.subheader("📉 Генератор презентаций")
    pres_topic = st.text_input("О чем презентация?")
    if st.button("Сгенерировать PPTX"):
        if pres_topic:
            with st.spinner("Готовлю слайды..."):
                text_content = ask_mistral(f"Напиши 3 ключевых пункта для презентации на тему: {pres_topic}")
                pptx_data = create_pptx(text_content)
                st.download_button("💾 Скачать презентацию", pptx_data, "presentation.pptx")
                st.success("Файл готов к скачиванию!")

# 6. Essay Checker
with tabs[5]:
    st.subheader("✍️ Проверка эссе")
    essay_text = st.text_area("Вставь текст эссе:", height=300)
    if st.button("Проверить"):
        if essay_text:
            with st.spinner("Проверяю грамотность и стиль..."):
                review = ask_mistral(f"Проверь это эссе на ошибки и дай рекомендации по стилю:\n{essay_text}")
                st.markdown(f'<div class="feature-card">{review}</div>', unsafe_allow_html=True)

# 7. Translator
with tabs[6]:
    st.subheader("🌍 Переводчик")
    col1, col2 = st.columns(2)
    with col1:
        text_to_translate = st.text_area("Текст для перевода:")
    with col2:
        target_lang = st.selectbox("На какой язык?", ["English", "German", "French", "Chinese", "Spanish"])
    if st.button("Перевести"):
        if text_to_translate:
            with st.spinner("Перевожу..."):
                translation = ask_mistral(f"Translate this text to {target_lang}:\n{text_to_translate}")
                st.markdown(f'<div class="feature-card">{translation}</div>', unsafe_allow_html=True)

# 8. Image Generator (Pollinations.ai)
with tabs[7]:
    st.subheader("🖼 Иллюстратор (Pollinations.ai)")
    img_prompt = st.text_input("Что нужно нарисовать?")
    if st.button("Создать изображение"):
        if img_prompt:
            with st.spinner("🎨 Рисую через Pollinations.ai..."):
                img_data = generate_image_pollinations(img_prompt)
                if img_data:
                    st.image(img_data, caption=f"Визуализация: {img_prompt}", width="stretch")
                    st.success("Изображение сгенерировано!")
